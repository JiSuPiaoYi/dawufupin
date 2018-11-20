from pdfminer.pdfparser import PDFParser, PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LTTextBoxHorizontal, LAParams
from pdfminer.pdfinterp import PDFTextExtractionNotAllowed
import docx
import xlrd
from pptx import Presentation

# 将txt文档的内容读取转化为str
def txt_parse(path):
    try:
        fileobj = open(path, 'r', encoding='utf-8')
        content = str(fileobj.read())
    except:
        try:
            fileobj = open(path, 'r', encoding='gbk')
            content = str(fileobj.read())
        except:
            try:
                fileobj = open(path, 'r', encoding='unicode')
                content = str(fileobj.read())
            except:
                fileobj = open(path, 'r')
                content = 'aaaaa'
    finally:
        fileobj.close()
        # print(content)
    return content



# 解析pdf文本，返回文本字符串
def pdf_parse(path):

    fp = open(path, 'rb')  # 以二进制读模式打开
    # 用文件对象来创建一个pdf文档分析器
    praser = PDFParser(fp)
    # 创建一个PDF文档
    doc = PDFDocument()
    # 连接分析器 与文档对象
    praser.set_document(doc)
    doc.set_parser(praser)

    # 提供初始化密码
    # 如果没有密码 就创建一个空的字符串
    doc.initialize()

    # 检测文档是否提供txt转换，不提供就忽略
    if not doc.is_extractable:
        return ''
        # raise PDFTextExtractionNotAllowed
    else:
        # 创建PDf 资源管理器 来管理共享资源
        rsrcmgr = PDFResourceManager()
        # 创建一个PDF设备对象
        laparams = LAParams()
        device = PDFPageAggregator(rsrcmgr, laparams=laparams)
        # 创建一个PDF解释器对象
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        result = ""
        # 循环遍历列表，每次处理一个page的内容
        for page in doc.get_pages():  # doc.get_pages() 获取page列表
            interpreter.process_page(page)
            # 接受该页面的LTPage对象
            layout = device.get_result()
            #  这里layout是一个LTPage对象 里面存放着 这个page解析出的各种对象
            # 一般包括LTTextBox, LTFigure, LTImage, LTTextBoxHorizontal 等等 想要获取文本就获得对象的text属性，
            for x in layout:
                if isinstance(x, LTTextBoxHorizontal):
                    result += x.get_text()

        return result


# 解析word文本，返回文本的字符串
def word_parse(path):
    doc = docx.Document(path)
    re = doc.paragraphs
    con = ""
    for p in re:
        con += p.text
    return con


# 解析excel文本，返回文本的字符串
def excel_parse(path):
    data = xlrd.open_workbook(path)
    tables = data.sheets()
    stri = ""
    for table in tables:
        nrows = table.nrows
        ncols = table.ncols
        for i in range(nrows):
            for j in range(ncols):
                stri += str(table.cell(i, j).value) + ","

    return stri


# 解析ppt，返回文本字符串
def ppt_parse2(path):

    prs = Presentation(path)
    text_runs = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs = text_runs + run.text + ','
    return text_runs

